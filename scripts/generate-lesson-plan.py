#!/usr/bin/env python3
"""
CTE Lesson Plan Generator

Generates:
1. Individual CTE format lesson plans (one per day)
2. Teacher handout document (all days in one document)
3. Student handouts when needed
4. Daily lesson presentations (PowerPoint with images and videos)

All files are organized into week-based folders (Week01, Week02, etc.)

Usage: python generate-lesson-plan.py '<json_data>'
"""

import sys
import json
import os
import requests
import tempfile
import re
from io import BytesIO
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from datetime import datetime

# PowerPoint imports for presentations
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# API KEYS AND CONFIGURATION
# ============================================================================

# Get Pexels API key from environment variable
# Set with: export PEXELS_API_KEY='your-api-key'
# Get a free key at: https://www.pexels.com/api/
PEXELS_API_KEY = os.environ.get('PEXELS_API_KEY', '')

# Preferred YouTube channels for video search
PREFERRED_YOUTUBE_CHANNELS = [
    'StudioBinder',
    'Film Riot',
    'D4Darious',
    'Cinecom',
    'Peter McKinnon',
    'Premiere Gal'
]

# Unit color themes (unit name -> (primary_color, secondary_color, accent_color))
UNIT_COLOR_THEMES = {
    'Introduction & History of Film': (PptxRGBColor(0x8B, 0x45, 0x13), PptxRGBColor(0xF5, 0xF5, 0xDC), PptxRGBColor(0xD2, 0x69, 0x1E)),  # Brown/sepia
    'Pre-Production': (PptxRGBColor(0x2E, 0x86, 0xAB), PptxRGBColor(0xE8, 0xF4, 0xF8), PptxRGBColor(0x56, 0xB4, 0xE9)),  # Blue
    'Camera Basics': (PptxRGBColor(0xE6, 0x55, 0x00), PptxRGBColor(0xFF, 0xF3, 0xE0), PptxRGBColor(0xFF, 0x8C, 0x00)),  # Orange
    'Premiere Pro Intro': (PptxRGBColor(0x9B, 0x59, 0xB6), PptxRGBColor(0xF5, 0xEE, 0xF8), PptxRGBColor(0xE9, 0x1E, 0x63)),  # Purple/Adobe
    'Advanced Techniques': (PptxRGBColor(0x1A, 0x1A, 0x2E), PptxRGBColor(0xE8, 0xE8, 0xE8), PptxRGBColor(0x00, 0xD4, 0xFF)),  # Dark/cyan
    'PSA Pre-Production': (PptxRGBColor(0x27, 0xAE, 0x60), PptxRGBColor(0xE8, 0xF8, 0xF0), PptxRGBColor(0x2E, 0xCC, 0x71)),  # Green
    'PSA Production': (PptxRGBColor(0x27, 0xAE, 0x60), PptxRGBColor(0xE8, 0xF8, 0xF0), PptxRGBColor(0x2E, 0xCC, 0x71)),  # Green
    'PSA Post-Production': (PptxRGBColor(0x27, 0xAE, 0x60), PptxRGBColor(0xE8, 0xF8, 0xF0), PptxRGBColor(0x2E, 0xCC, 0x71)),  # Green
    'News Segment': (PptxRGBColor(0xC0, 0x39, 0x2B), PptxRGBColor(0xFD, 0xED, 0xEC), PptxRGBColor(0xE7, 0x4C, 0x3C)),  # Red/news
    'News/Documentary Intro': (PptxRGBColor(0x5D, 0x4E, 0x37), PptxRGBColor(0xF5, 0xF0, 0xE6), PptxRGBColor(0x8B, 0x7D, 0x6B)),  # Earth tones
    'Documentary Production': (PptxRGBColor(0x5D, 0x4E, 0x37), PptxRGBColor(0xF5, 0xF0, 0xE6), PptxRGBColor(0x8B, 0x7D, 0x6B)),  # Earth tones
    'Documentary Post': (PptxRGBColor(0x5D, 0x4E, 0x37), PptxRGBColor(0xF5, 0xF0, 0xE6), PptxRGBColor(0x8B, 0x7D, 0x6B)),  # Earth tones
    'Music Video Pre-Production': (PptxRGBColor(0xE9, 0x1E, 0x63), PptxRGBColor(0xFC, 0xE4, 0xEC), PptxRGBColor(0x9C, 0x27, 0xB0)),  # Pink/purple
    'Music Video Production': (PptxRGBColor(0xE9, 0x1E, 0x63), PptxRGBColor(0xFC, 0xE4, 0xEC), PptxRGBColor(0x9C, 0x27, 0xB0)),  # Pink/purple
    'Music Video Post': (PptxRGBColor(0xE9, 0x1E, 0x63), PptxRGBColor(0xFC, 0xE4, 0xEC), PptxRGBColor(0x9C, 0x27, 0xB0)),  # Pink/purple
    'Final Exam': (PptxRGBColor(0x1a, 0x3c, 0x6e), PptxRGBColor(0xD6, 0xE3, 0xF8), PptxRGBColor(0x34, 0x98, 0xDB)),  # Navy (default)
}

# Default theme (navy)
DEFAULT_COLOR_THEME = (PptxRGBColor(0x1a, 0x3c, 0x6e), PptxRGBColor(0xD6, 0xE3, 0xF8), PptxRGBColor(0x34, 0x98, 0xDB))

# Shared document color palette (used across teacher/student handouts)
DOC_NAVY_BLUE = RGBColor(0x1a, 0x3c, 0x6e)
DOC_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
DOC_MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
DOC_LIGHT_BLUE = "D6E3F8"
DOC_LIGHT_GRAY = "F5F5F5"
DOC_ACCENT_BLUE = "4A90D9"
DOC_CREAM_YELLOW = "FFF9E6"
DOC_SOFT_GREEN = "E8F5E9"


# ============================================================================
# PEXELS API FUNCTIONS
# ============================================================================

def search_pexels_image(query, per_page=1):
    """Search Pexels for an image matching the query. Returns image URL or None."""
    if not PEXELS_API_KEY:
        return None
    try:
        headers = {'Authorization': PEXELS_API_KEY}
        params = {'query': query, 'per_page': per_page, 'orientation': 'landscape'}
        response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params, timeout=10)

        if response.status_code == 200:
            data = response.json()
            if data.get('photos') and len(data['photos']) > 0:
                return data['photos'][0]['src']['large']
        elif response.status_code == 401:
            print("Pexels API error: Invalid API key", file=sys.stderr)
        return None
    except requests.exceptions.Timeout:
        print("Pexels API error: Request timed out", file=sys.stderr)
        return None
    except requests.exceptions.ConnectionError:
        print("Pexels API error: Could not connect", file=sys.stderr)
        return None
    except Exception as e:
        print(f"Pexels API error: {e}", file=sys.stderr)
        return None


def download_image(url):
    """Download an image from URL and return as BytesIO object."""
    if not url or not url.startswith('https://'):
        return None
    try:
        response = requests.get(url, timeout=15)
        if response.status_code == 200:
            content_type = response.headers.get('Content-Type', '')
            if not content_type.startswith('image/'):
                return None
            return BytesIO(response.content)
        return None
    except requests.exceptions.Timeout:
        print("Image download error: Request timed out", file=sys.stderr)
        return None
    except requests.exceptions.ConnectionError:
        print("Image download error: Could not connect", file=sys.stderr)
        return None
    except Exception as e:
        print(f"Image download error: {e}", file=sys.stderr)
        return None


def get_topic_image(topic, context="media production"):
    """Get an image for a topic. Returns BytesIO image data or None."""
    # Build search query with context
    search_terms = [
        f"{topic} {context}",
        f"{topic} film",
        f"{topic} video",
        topic
    ]

    for query in search_terms:
        image_url = search_pexels_image(query)
        if image_url:
            image_data = download_image(image_url)
            if image_data:
                return image_data, image_url

    return None, None


# ============================================================================
# YOUTUBE VIDEO SEARCH FUNCTIONS
# ============================================================================

# Curated video library for common media production topics
# Format: keyword -> (video_url, video_title)
CURATED_VIDEOS = {
    'camera angles': ('https://www.youtube.com/watch?v=SlNviMsi0K0', 'Camera Angles Explained - StudioBinder'),
    'shot types': ('https://www.youtube.com/watch?v=AyML8xuKfoc', 'Ultimate Guide to Camera Shots - StudioBinder'),
    'composition': ('https://www.youtube.com/watch?v=O8i7OKbWmRM', 'Composition in Film - StudioBinder'),
    'rule of thirds': ('https://www.youtube.com/watch?v=O8i7OKbWmRM', 'Composition Techniques - StudioBinder'),
    'aperture': ('https://www.youtube.com/watch?v=YojL7UQTVhc', 'Aperture Explained - Film Riot'),
    'shutter speed': ('https://www.youtube.com/watch?v=HYB67U89sKs', 'Shutter Speed for Video'),
    'iso': ('https://www.youtube.com/watch?v=WEApFs0aDaE', 'ISO Explained for Filmmakers'),
    'exposure triangle': ('https://www.youtube.com/watch?v=3eVjUrY9a9c', 'Exposure Triangle - Peter McKinnon'),
    'lighting': ('https://www.youtube.com/watch?v=j_Sov3xmgwg', 'Cinematic Lighting Techniques'),
    'three point lighting': ('https://www.youtube.com/watch?v=j_Sov3xmgwg', '3 Point Lighting Setup'),
    'audio': ('https://www.youtube.com/watch?v=U1dMlVwUsrA', 'Audio Recording for Film'),
    'premiere pro': ('https://www.youtube.com/watch?v=Hls3Tp7JS8E', 'Premiere Pro Tutorial for Beginners'),
    'editing': ('https://www.youtube.com/watch?v=O6ERELse_QY', 'Video Editing Basics - Film Riot'),
    'color correction': ('https://www.youtube.com/watch?v=lxHnCXZgeQc', 'Color Correction in Premiere Pro'),
    'color grading': ('https://www.youtube.com/watch?v=lxHnCXZgeQc', 'Color Grading Tutorial'),
    'copyright': ('https://www.youtube.com/watch?v=1Jwo5qc78QU', 'Copyright for Filmmakers'),
    'fair use': ('https://www.youtube.com/watch?v=1Jwo5qc78QU', 'Fair Use Explained'),
    'storyboard': ('https://www.youtube.com/watch?v=RQsvhq28sOI', 'How to Storyboard - StudioBinder'),
    'screenplay': ('https://www.youtube.com/watch?v=vrvawtrRxsw', 'Screenwriting Basics'),
    'script': ('https://www.youtube.com/watch?v=vrvawtrRxsw', 'How to Write a Script'),
    'documentary': ('https://www.youtube.com/watch?v=fMF0xQo-E8U', 'Documentary Filmmaking Tips'),
    'interview': ('https://www.youtube.com/watch?v=R0LD7VHxYiE', 'How to Film an Interview'),
    'b-roll': ('https://www.youtube.com/watch?v=mHZ6LGKnDc0', 'B-Roll Techniques - Film Riot'),
    'green screen': ('https://www.youtube.com/watch?v=hRsrVjbYyiE', 'Green Screen Tutorial'),
    'foley': ('https://www.youtube.com/watch?v=U_tqB4IZvMk', 'Foley Sound Effects Explained'),
    'sound design': ('https://www.youtube.com/watch?v=U_tqB4IZvMk', 'Sound Design for Film'),
    'music video': ('https://www.youtube.com/watch?v=p5rQHoaQpTw', 'How to Make a Music Video'),
    'psa': ('https://www.youtube.com/watch?v=9sjkvYdoH9o', 'How to Make a PSA'),
    'news': ('https://www.youtube.com/watch?v=vMnTZrFa-Wc', 'Broadcast News Production'),
    'film history': ('https://www.youtube.com/watch?v=HCYJBwY-Qsc', 'History of Cinema'),
}


def search_youtube_video(topic, preferred_channels=None):
    """
    Find a YouTube video for a topic.
    First checks curated library, then attempts web search.
    Returns (video_url, video_title) or (None, None).
    """
    # Normalize topic for matching
    topic_lower = topic.lower().strip()

    # Check curated videos first (most reliable)
    for keyword, (url, title) in CURATED_VIDEOS.items():
        if keyword in topic_lower or topic_lower in keyword:
            return url, title

    # Check for partial matches
    topic_words = topic_lower.split()
    for keyword, (url, title) in CURATED_VIDEOS.items():
        for word in topic_words:
            if len(word) > 3 and word in keyword:
                return url, title

    # Fallback: Try web search
    try:
        import warnings
        warnings.filterwarnings('ignore')
        from duckduckgo_search import DDGS

        search_queries = [
            f'{topic} filmmaking tutorial youtube',
            f'{topic} video production tutorial',
        ]

        with DDGS() as ddgs:
            for query in search_queries:
                try:
                    results = list(ddgs.text(query, max_results=10))
                    for result in results:
                        url = result.get('href', '')
                        title = result.get('title', '')
                        if 'youtube.com/watch' in url or 'youtu.be/' in url:
                            return url, title
                except Exception:
                    continue
    except Exception:
        pass

    return None, None


def get_youtube_video_id(url):
    """Extract YouTube video ID from URL."""
    if not url:
        return None

    # Handle youtu.be format
    if 'youtu.be/' in url:
        match = re.search(r'youtu\.be/([a-zA-Z0-9_-]+)', url)
        if match:
            return match.group(1)

    # Handle youtube.com/watch format
    if 'youtube.com/watch' in url:
        match = re.search(r'[?&]v=([a-zA-Z0-9_-]+)', url)
        if match:
            return match.group(1)

    return None


# Path to CTE lesson plan template (Word document)
# Set with: export CTE_TEMPLATE_PATH='/path/to/template.docx'
# Or place template in ./templates/ directory
TEMPLATE_PATH = os.environ.get('CTE_TEMPLATE_PATH',
    os.path.join(os.path.dirname(os.path.dirname(__file__)), 'templates', 'CTE_Lesson_Plan_Template.docx'))

# Output directory for generated files
# Set with: export CTE_OUTPUT_DIR='/path/to/output'
# Defaults to ./output/ in the project directory
OUTPUT_DIR = os.environ.get('CTE_OUTPUT_DIR',
    os.path.join(os.path.dirname(os.path.dirname(__file__)), 'output'))

# Checkbox mappings
MATERIALS_CHECKBOXES = {
    'textbook': 'Textbook',
    'lab_manual': 'Lab Manual',
    'video_dvd': 'Video/DVD',
    'labs': 'Labs',
    'posters': 'Posters',
    'speaker': 'Speaker',
    'projector': 'Projector',
    'computer': 'Computer',
    'supplemental_materials': 'Supplemental Materials',
    'student_journals': 'Student Journals',
    'other_equipment': 'Other Equipment'
}

METHODS_CHECKBOXES = {
    'discussion': 'Discussion',
    'demonstration': 'Demonstration',
    'lecture': 'Lecture',
    'powerpoint': 'Power Point',
    'multimedia': 'Multi-Media',
    'guest_speaker': 'Guest Speaker'
}

ASSESSMENT_CHECKBOXES = {
    'homework': 'Homework',
    'classwork': 'Classwork',
    'test': 'Test',
    'project_based': 'Project-based',
    'teamwork': 'Teamwork',
    'observation': 'Teacher Observation',
    'performance': 'Performance',
    'on_task': 'On-Task',
    'other': 'Other'
}

CURRICULUM_CHECKBOXES = {
    'math': 'Math',
    'science': 'Science',
    'reading': 'Reading',
    'social_studies': 'Social Studies',
    'english': 'English',
    'government_economics': 'Government/Economics',
    'fine_arts': 'Fine Arts',
    'foreign_language': 'Foreign Language',
    'technology': 'Technology'
}

OTHER_AREAS_CHECKBOXES = {
    'safety': 'Safety',
    'management_skills': 'Management Skills',
    'teamwork': 'Teamwork',
    'live_work': 'Live work',
    'higher_order_reasoning': 'Higher Order Reasoning',
    'varied_learning': 'Varied Learning',
    'work_ethics': 'Work Ethics',
    'integrated_academics': 'Integrated Academics',
    'ctso': 'CTSO',
    'problem_solving': 'Problem Solving'
}


def setup_document(doc, top=0.6, bottom=0.6, left=0.7, right=0.7, line_spacing=1.15, space_after=6):
    """Configure default document styles and margins."""
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = DOC_DARK_GRAY
    style.paragraph_format.space_after = Pt(space_after)
    style.paragraph_format.line_spacing = line_spacing

    for section in doc.sections:
        section.top_margin = Inches(top)
        section.bottom_margin = Inches(bottom)
        section.left_margin = Inches(left)
        section.right_margin = Inches(right)


def add_doc_section_header(doc, text, level=1, accent_color=None):
    """Add a section header with left accent bar to a document."""
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    color = accent_color or DOC_ACCENT_BLUE
    header_tbl = doc.add_table(rows=1, cols=2)
    header_tbl.autofit = False

    sidebar_cell = header_tbl.rows[0].cells[0]
    sidebar_cell.width = Inches(0.08)
    sidebar_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}" w:val="clear"/>')
    sidebar_cell._tc.get_or_add_tcPr().append(sidebar_shading)

    content_cell = header_tbl.rows[0].cells[1]
    content_cell.width = Inches(6.8)

    p = content_cell.paragraphs[0]
    run = p.add_run(text)
    run.bold = True
    run.font.color.rgb = DOC_NAVY_BLUE
    if level == 1:
        run.font.size = Pt(16)
        run.font.name = 'Cambria'
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(4)
    else:
        run.font.size = Pt(13)
        run.font.name = 'Cambria'
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after = Pt(2)

    return p


def add_doc_header_banner(doc, title_text, subtitle_text=None, badge_text=None):
    """Add a styled header banner with accent bar to a document."""
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    header_table = doc.add_table(rows=2, cols=1)
    header_table.style = 'Table Grid'

    # Accent bar
    accent_cell = header_table.rows[0].cells[0]
    accent_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{DOC_ACCENT_BLUE}" w:val="clear"/>')
    accent_cell._tc.get_or_add_tcPr().append(accent_shading)
    accent_p = accent_cell.paragraphs[0]
    accent_p.paragraph_format.space_after = Pt(0)
    tr = header_table.rows[0]._tr
    trPr = tr.get_or_add_trPr()
    trHeight = parse_xml(f'<w:trHeight {nsdecls("w")} w:val="100" w:hRule="exact"/>')
    trPr.append(trHeight)

    # Main header cell
    main_cell = header_table.rows[1].cells[0]
    main_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
    main_cell._tc.get_or_add_tcPr().append(main_shading)

    # Optional badge (e.g., "WEEK 5")
    if badge_text:
        p = main_cell.paragraphs[0]
        run = p.add_run(badge_text)
        run.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xD6, 0xE3, 0xF8)
        run.font.name = 'Calibri'
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(2)
        # Title goes in next paragraph
        p = main_cell.add_paragraph()
    else:
        p = main_cell.paragraphs[0]

    # Title
    run = p.add_run(title_text)
    run.bold = True
    run.font.size = Pt(24) if not badge_text else Pt(28)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = 'Cambria'
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(4 if badge_text else 8)
    p.paragraph_format.space_after = Pt(4)

    # Subtitle
    if subtitle_text:
        p = main_cell.add_paragraph()
        run = p.add_run(subtitle_text)
        run.font.size = Pt(11 if badge_text else 12)
        run.font.color.rgb = RGBColor(0xD6, 0xE3, 0xF8)
        run.font.name = 'Calibri'
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(8)

    return main_cell


def sanitize_slug(text, max_length=25):
    """Sanitize text for use in filenames, removing unsafe characters."""
    slug = re.sub(r'[^\w\s-]', '', text)
    slug = slug.replace(' ', '_')
    slug = re.sub(r'[_.]+', lambda m: m.group(0)[0], slug)
    return slug[:max_length].rstrip('_')


def get_week_folder(week_num):
    """Get the week folder path, creating it if needed."""
    week_str = str(int(week_num)).zfill(2)
    week_folder = os.path.join(OUTPUT_DIR, f"Week{week_str}")
    os.makedirs(week_folder, exist_ok=True)
    return week_folder


def remove_red_text(doc):
    """Remove red color from all text in the document, making it black."""
    BLACK = RGBColor(0, 0, 0)

    # Process all tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        # Check if text has red color
                        if run.font.color.rgb:
                            color = run.font.color.rgb
                            # If it's red-ish, change to black
                            if color[0] > 150 and color[1] < 100 and color[2] < 100:
                                run.font.color.rgb = BLACK
                        # Also ensure any new text is black
                        if run.font.color.type is None:
                            run.font.color.rgb = BLACK

    # Process paragraphs outside tables
    for para in doc.paragraphs:
        for run in para.runs:
            if run.font.color.rgb:
                color = run.font.color.rgb
                if color[0] > 150 and color[1] < 100 and color[2] < 100:
                    run.font.color.rgb = BLACK


def mark_checkboxes_in_cell(cell, checkbox_map, selected_items):
    """Mark checkboxes in a cell by replacing underscores with checkmarks."""
    import re
    for para in cell.paragraphs:
        for run in para.runs:
            text = run.text
            for key, label in checkbox_map.items():
                if key in selected_items:
                    pattern = r'_+(\s*)' + re.escape(label)
                    if re.search(pattern, text, re.IGNORECASE):
                        text = re.sub(pattern, r'X\1' + label, text, flags=re.IGNORECASE)
            run.text = text


def set_cell_text(cell, text):
    """Set the text of a cell, preserving the first paragraph's formatting but ensuring black text."""
    BLACK = RGBColor(0, 0, 0)

    while len(cell.paragraphs) > 1:
        p = cell.paragraphs[-1]
        p._element.getparent().remove(p._element)

    if cell.paragraphs:
        para = cell.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
            para.runs[0].font.color.rgb = BLACK
        else:
            run = para.add_run(text)
            run.font.color.rgb = BLACK
    else:
        para = cell.add_paragraph()
        run = para.add_run(text)
        run.font.color.rgb = BLACK


def build_procedures_text(day_data):
    """Build the Procedures/Activities/Learning Experiences text from schedule data."""
    procedures = day_data.get('procedures', '')
    if procedures:
        return procedures

    # Build from schedule if no explicit procedures provided
    schedule = day_data.get('schedule', [])
    if not schedule:
        return ''

    lines = []
    for activity in schedule:
        if isinstance(activity, dict):
            time = activity.get('time', activity.get('duration', ''))
            name = activity.get('name', activity.get('activity', ''))
            desc = activity.get('description', '')
            if time and name:
                lines.append(f"{time} - {name}: {desc}" if desc else f"{time} - {name}")
            elif name:
                lines.append(f"{name}: {desc}" if desc else name)
        else:
            lines.append(str(activity))

    return '\n'.join(lines)


def build_differentiation_text(day_data):
    """Build the Provision for Individual Differences text from differentiation data."""
    individual_diff = day_data.get('individual_differences', '')
    if individual_diff:
        return individual_diff

    # Build from differentiation if no explicit text provided
    diff = day_data.get('differentiation', {})
    if not diff:
        return ''

    if isinstance(diff, str):
        return diff

    lines = []
    if diff.get('Advanced'):
        lines.append(f"Advanced Learners: {diff['Advanced']}")
    if diff.get('Struggling'):
        lines.append(f"Struggling Learners: {diff['Struggling']}")
    if diff.get('ELL'):
        lines.append(f"ELL Students: {diff['ELL']}")

    # Handle any other differentiation levels
    for level, strategy in diff.items():
        if level not in ['Advanced', 'Struggling', 'ELL']:
            lines.append(f"{level}: {strategy}")

    return '\n'.join(lines)


def get_lesson_text(day_data, include_materials=False):
    """Extract all searchable text content from a day's lesson data."""
    parts = [
        day_data.get('topic', ''),
        day_data.get('overview', ''),
        ' '.join(day_data.get('objectives', [])),
    ]
    if include_materials:
        parts.append(' '.join(day_data.get('day_materials', [])))
    for activity in day_data.get('schedule', []):
        if isinstance(activity, dict):
            parts.append(activity.get('name', ''))
            parts.append(activity.get('description', ''))
    return ' '.join(parts).lower()


def check_keywords(text, keywords):
    """Check if any keyword is found in the given text."""
    return any(kw in text for kw in keywords)


def infer_other_areas(day_data, curriculum_areas):
    """Infer other areas addressed based on lesson content."""
    other_areas = list(day_data.get('other_areas', []))
    all_text = get_lesson_text(day_data)

    keyword_map = {
        'safety': ['safety', 'equipment', 'handling', 'protective', 'hazard', 'proper use', 'safely', 'precaution'],
        'management_skills': ['time management', 'organize', 'planning', 'schedule', 'project management', 'workflow', 'deadline'],
        'teamwork': ['team', 'group', 'collaborat', 'partner', 'cooperative', 'crew', 'together'],
        'live_work': ['client', 'real-world', 'live production', 'actual client', 'community partner'],
        'higher_order_reasoning': ['analyze', 'evaluat', 'create', 'critiqu', 'compare', 'synthesize', 'design', 'develop', 'assess'],
        'work_ethics': ['professional', 'responsibility', 'deadline', 'punctual', 'quality', 'ethic', 'industry standard'],
        'ctso': ['skillsusa', 'ctso', 'competition', 'career development', 'leadership'],
        'problem_solving': ['problem', 'solve', 'troubleshoot', 'debug', 'fix', 'challenge', 'solution', 'figure out'],
    }

    for area, keywords in keyword_map.items():
        if area not in other_areas and check_keywords(all_text, keywords):
            other_areas.append(area)

    # Varied Learning - also checks differentiation data
    if 'varied_learning' not in other_areas:
        if check_keywords(all_text, ['visual', 'hands-on', 'demonstration', 'practice', 'kinesthetic', 'auditory']) or \
           bool(day_data.get('differentiation')):
            other_areas.append('varied_learning')

    # Integrated Academics - if curriculum areas are checked
    if curriculum_areas and 'integrated_academics' not in other_areas:
        other_areas.append('integrated_academics')

    return other_areas


def infer_curriculum_areas(day_data):
    """Infer integrated curriculum areas based on lesson content."""
    curriculum = list(day_data.get('curriculum', []))
    all_text = get_lesson_text(day_data)

    keyword_map = {
        'technology': ['camera', 'editing', 'software', 'premiere', 'photoshop', 'computer', 'digital', 'video', 'audio', 'equipment'],
        'english': ['script', 'writing', 'story', 'narrative', 'reading', 'research', 'interview', 'article', 'news'],
        'fine_arts': ['composition', 'visual', 'design', 'aesthetic', 'creative', 'artistic', 'color', 'lighting', 'framing'],
        'math': ['exposure', 'ratio', 'frame rate', 'aperture', 'shutter speed', 'iso', 'calculation', 'percentage'],
        'science': ['light', 'sound wave', 'physics', 'optics', 'frequency', 'wavelength'],
        'social_studies': ['history', 'documentary', 'social', 'community', 'culture', 'news', 'current events', 'psa', 'public service'],
    }

    for area, keywords in keyword_map.items():
        if area not in curriculum and check_keywords(all_text, keywords):
            curriculum.append(area)

    # Reading is a sub-check of english
    if 'reading' not in curriculum and check_keywords(all_text, ['reading', 'research', 'article']):
        curriculum.append('reading')

    return curriculum


def infer_materials(day_data):
    """Infer materials and equipment based on lesson content."""
    materials = list(day_data.get('materials', []))
    all_text = get_lesson_text(day_data, include_materials=True)

    keyword_map = {
        'projector': ['presentation', 'present', 'show', 'display', 'screen', 'projector', 'slides', 'powerpoint'],
        'computer': ['computer', 'premiere', 'photoshop', 'editing', 'software', 'digital', 'laptop', 'workstation'],
        'video_dvd': ['video', 'watch', 'film', 'movie', 'clip', 'example', 'youtube', 'dvd'],
        'labs': ['lab', 'studio', 'hands-on', 'practice', 'filming', 'shoot', 'record'],
        'speaker': ['audio', 'sound', 'music', 'listen', 'speaker', 'playback'],
        'supplemental_materials': ['handout', 'worksheet', 'guide', 'reference', 'template', 'storyboard', 'script'],
        'other_equipment': ['camera', 'tripod', 'lighting', 'light', 'microphone', 'mic', 'equipment', 'gear', 'sd card', 'memory card'],
        'student_journals': ['journal', 'notebook', 'notes', 'reflection', 'write', 'record thoughts'],
        'posters': ['poster', 'chart', 'diagram', 'visual aid', 'infographic'],
    }

    for material, keywords in keyword_map.items():
        if material not in materials and check_keywords(all_text, keywords):
            materials.append(material)

    return materials


def infer_methods(day_data):
    """Infer instructional methods based on lesson content."""
    methods = list(day_data.get('methods', []))
    all_text = get_lesson_text(day_data)

    # Extract activity names for special lecture check
    activity_names = []
    for activity in day_data.get('schedule', []):
        if isinstance(activity, dict):
            activity_names.append(activity.get('name', '').lower())

    keyword_map = {
        'discussion': ['discussion', 'discuss', 'debate', 'share', 'q&a', 'conversation', 'talk about'],
        'demonstration': ['demonstrat', 'show how', 'model', 'walk through', 'example', 'tutorial'],
        'powerpoint': ['powerpoint', 'presentation', 'slides', 'slide deck', 'pptx'],
        'multimedia': ['video', 'multimedia', 'multi-media', 'youtube', 'film', 'audio', 'digital'],
        'guest_speaker': ['guest speaker', 'guest', 'industry professional', 'visitor', 'expert'],
    }

    for method, keywords in keyword_map.items():
        if method not in methods and check_keywords(all_text, keywords):
            methods.append(method)

    # Lecture - also checks activity names directly
    if 'lecture' not in methods:
        if any(name in ['direct instruction', 'lecture', 'mini-lecture', 'instruction'] for name in activity_names) or \
           check_keywords(all_text, ['lecture', 'direct instruction', 'teach', 'explain', 'present content', 'introduce']):
            methods.append('lecture')

    return methods


def infer_assessment(day_data):
    """Infer assessment strategies based on lesson content."""
    assessment = list(day_data.get('assessment', []))
    all_text = get_lesson_text(day_data)

    keyword_map = {
        'classwork': ['classwork', 'class work', 'activity', 'practice', 'exercise', 'in-class', 'work on'],
        'observation': ['observ', 'monitor', 'circulate', 'watch', 'check in', 'walk around'],
        'project_based': ['project', 'final', 'deliverable', 'portfolio', 'create', 'produce', 'video project'],
        'teamwork': ['team', 'group', 'partner', 'collaborat', 'crew', 'together', 'peer'],
        'performance': ['perform', 'present', 'demonstrat', 'show', 'pitch', 'share out'],
        'on_task': ['participat', 'engag', 'on-task', 'focused', 'active'],
        'test': ['test', 'quiz', 'exam', 'assessment'],
        'homework': ['homework', 'home work', 'take home', 'assignment', 'due next'],
    }

    for strategy, keywords in keyword_map.items():
        if strategy not in assessment and check_keywords(all_text, keywords):
            assessment.append(strategy)

    # Exit ticket also counts as classwork
    if 'classwork' not in assessment and check_keywords(all_text, ['exit ticket', 'exit slip', 'reflection']):
        assessment.append('classwork')

    return assessment


def build_overview_text(day_data):
    """Build overview text from lesson data if not explicitly provided."""
    # If overview is explicitly provided, use it
    if day_data.get('overview'):
        return day_data['overview']

    # Otherwise, build from available data
    topic = day_data.get('topic', 'the lesson topic')
    objectives = day_data.get('objectives', [])

    # Start with the topic
    overview_parts = [f"Students will learn about {topic}."]

    # Add objectives summary
    if objectives:
        if len(objectives) == 1:
            overview_parts.append(f"The primary objective is to {objectives[0].lower().lstrip('students will ').lstrip('to ')}.")
        else:
            overview_parts.append(f"Key objectives include: {objectives[0].lower().lstrip('students will ').lstrip('to ')}")
            for obj in objectives[1:]:
                clean_obj = obj.lower().lstrip('students will ').lstrip('to ')
                overview_parts.append(f"and {clean_obj}")

    # Add activity highlights from schedule
    schedule = day_data.get('schedule', [])
    activities = []
    for activity in schedule:
        if isinstance(activity, dict):
            name = activity.get('name', '').lower()
            if 'hands-on' in name or 'activity' in name or 'practice' in name:
                activities.append(activity.get('description', ''))

    if activities:
        overview_parts.append(f"Students will engage in hands-on activities including: {activities[0][:100]}...")

    return ' '.join(overview_parts)


def generate_cte_lesson_plan(day_data, week_num, day_num):
    """Generate a single CTE format lesson plan document."""
    doc = Document(TEMPLATE_PATH)
    table = doc.tables[0]

    # Build auto-generated fields
    procedures_text = build_procedures_text(day_data)
    differentiation_text = build_differentiation_text(day_data)
    overview_text = build_overview_text(day_data)
    curriculum_areas = infer_curriculum_areas(day_data)
    other_areas = infer_other_areas(day_data, curriculum_areas)
    materials = infer_materials(day_data)
    methods = infer_methods(day_data)
    assessment = infer_assessment(day_data)

    # Fill in all fields
    set_cell_text(table.rows[1].cells[0], f"Week: {week_num}")
    set_cell_text(table.rows[1].cells[1], f"Course Title: Media Foundations")
    set_cell_text(table.rows[2].cells[0], f"Topic: {day_data.get('topic', '')}")
    set_cell_text(table.rows[2].cells[1], f"Estimate duration in minutes: {day_data.get('duration', '90')}")
    set_cell_text(table.rows[5].cells[0], day_data.get('content_standards', ''))
    set_cell_text(table.rows[7].cells[0], overview_text)
    mark_checkboxes_in_cell(table.rows[7].cells[1], MATERIALS_CHECKBOXES, materials)
    set_cell_text(table.rows[9].cells[0], procedures_text)
    mark_checkboxes_in_cell(table.rows[11].cells[0], METHODS_CHECKBOXES, methods)
    mark_checkboxes_in_cell(table.rows[13].cells[0], ASSESSMENT_CHECKBOXES, assessment)
    set_cell_text(table.rows[13].cells[2], differentiation_text)
    mark_checkboxes_in_cell(table.rows[15].cells[0], CURRICULUM_CHECKBOXES, curriculum_areas)
    set_cell_text(table.rows[15].cells[2], day_data.get('embedded_credit', ''))
    mark_checkboxes_in_cell(table.rows[17].cells[0], OTHER_AREAS_CHECKBOXES, other_areas)
    set_cell_text(table.rows[17].cells[2], day_data.get('lesson_evaluation', ''))

    # Remove all red text from the document
    remove_red_text(doc)

    # Generate filename in week folder
    week_folder = get_week_folder(week_num)
    topic_slug = sanitize_slug(day_data.get('topic', 'Lesson'))
    filename = f"Day{day_num}_{topic_slug}_CTE.docx"
    output_path = os.path.join(week_folder, filename)

    doc.save(output_path)
    return output_path


def generate_teacher_handout(week_data):
    """Generate a professionally styled Canva-quality teacher handout."""
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    doc = Document()

    # Use shared color palette
    NAVY_BLUE = DOC_NAVY_BLUE
    DARK_GRAY = DOC_DARK_GRAY
    MEDIUM_GRAY = DOC_MEDIUM_GRAY
    LIGHT_BLUE = DOC_LIGHT_BLUE
    LIGHT_GRAY = DOC_LIGHT_GRAY
    ACCENT_BLUE = DOC_ACCENT_BLUE
    CREAM_YELLOW = DOC_CREAM_YELLOW
    SOFT_GREEN = DOC_SOFT_GREEN

    setup_document(doc)

    week_num = week_data.get('week', '')
    unit_name = week_data.get('unit', '')

    # === HEADER BANNER ===
    add_doc_header_banner(doc, unit_name,
                          subtitle_text="Media Foundations · Teacher Guide",
                          badge_text=f"WEEK {week_num}")

    doc.add_paragraph()

    def add_section_header(text, level=1):
        return add_doc_section_header(doc, text, level=level)

    # === Helper function for card-style boxes ===
    def add_card_box(content, bg_color=LIGHT_BLUE, border_color="1A3C6E"):
        """Create a card-style box with colored background."""
        card_table = doc.add_table(rows=1, cols=1)
        card_table.style = 'Table Grid'
        cell = card_table.rows[0].cells[0]

        # Background color
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)

        # Add content
        p = cell.paragraphs[0]
        if isinstance(content, str):
            run = p.add_run(content)
            run.font.color.rgb = DARK_GRAY

        return cell

    # === Helper function for pull-quote/tip boxes ===
    def add_tip_box(content, tip_type="tip"):
        """Create a styled tip/note box."""
        tip_colors = {'tip': SOFT_GREEN, 'warning': "FFF3CD", 'note': CREAM_YELLOW, 'important': "FFEBEE"}

        box_color = tip_colors.get(tip_type, CREAM_YELLOW)

        tip_table = doc.add_table(rows=1, cols=1)
        tip_table.style = 'Table Grid'
        cell = tip_table.rows[0].cells[0]

        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{box_color}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)

        p = cell.paragraphs[0]
        run = p.add_run(content)
        run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(10)

        return cell

    # === Helper function for styled tables ===
    def style_table(table, header_color=LIGHT_BLUE):
        # Set table width to full page
        table.autofit = True

        # Style header row
        if table.rows:
            for cell in table.rows[0].cells:
                # Set background color for header
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{header_color}" w:val="clear"/>')
                cell._tc.get_or_add_tcPr().append(shading)
                # Style header text
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(11)
                        run.font.color.rgb = NAVY_BLUE

    # === WEEK OVERVIEW BOX (Card style with sidebar accent) ===
    if week_data.get('week_overview') or week_data.get('week_focus'):
        add_section_header("Week Overview", level=1)

        overview_table = doc.add_table(rows=1, cols=1)
        overview_table.style = 'Table Grid'

        # Style the overview box with light blue background
        cell = overview_table.rows[0].cells[0]
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_BLUE}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)

        # Focus line with emphasis
        if week_data.get('week_focus'):
            p = cell.paragraphs[0]
            run = p.add_run("Focus: ")
            run.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = NAVY_BLUE
            focus_run = p.add_run(week_data['week_focus'])
            focus_run.font.color.rgb = DARK_GRAY
            focus_run.font.size = Pt(11)
            p.paragraph_format.space_after = Pt(8)

        # Overview content
        if week_data.get('week_overview'):
            p = cell.add_paragraph() if week_data.get('week_focus') else cell.paragraphs[0]
            run = p.add_run(week_data['week_overview'])
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(11)

        doc.add_paragraph()

    # === WEEKLY LEARNING OBJECTIVES ===
    if week_data.get('week_objectives'):
        add_section_header("Weekly Learning Objectives", level=1)

        # Create a styled card for objectives
        obj_table = doc.add_table(rows=0, cols=2)
        obj_table.style = 'Table Grid'
        obj_table.autofit = False

        for idx, obj in enumerate(week_data['week_objectives']):
            row = obj_table.add_row()

            # Number badge cell
            num_cell = row.cells[0]
            num_cell.width = Inches(0.4)
            p = num_cell.paragraphs[0]
            run = p.add_run(str(idx + 1))
            run.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(11)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Navy circle background
            num_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
            num_cell._tc.get_or_add_tcPr().append(num_shading)

            # Objective text cell
            text_cell = row.cells[1]
            p = text_cell.paragraphs[0]
            run = p.add_run(obj)
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(11)

        doc.add_paragraph()

    # === MATERIALS NEEDED FOR THE WEEK ===
    if week_data.get('week_materials'):
        add_section_header("Materials Needed for the Week", level=1)

        # Two-column layout for materials
        mat_table = doc.add_table(rows=0, cols=2)
        mat_table.style = 'Table Grid'
        mat_table.autofit = False

        materials = week_data['week_materials']
        for i in range(0, len(materials), 2):
            row = mat_table.add_row()
            for j in range(2):
                if i + j < len(materials):
                    cell = row.cells[j]
                    cell.width = Inches(3.4)
                    p = cell.paragraphs[0]
                    run = p.add_run(f"[ ] {materials[i + j]}")
                    run.font.color.rgb = DARK_GRAY
                    run.font.size = Pt(11)
                    # Alternating row color
                    if (i // 2) % 2 == 1:
                        mat_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_GRAY}" w:val="clear"/>')
                        cell._tc.get_or_add_tcPr().append(mat_shading)

        doc.add_paragraph()

    # === ASSESSMENT OVERVIEW ===
    if week_data.get('assessment_overview') or week_data.get('formative_assessment') or week_data.get('summative_assessment') or week_data.get('weekly_deliverable'):
        add_section_header("Assessment Overview", level=1)

        # Create assessment cards in a 3-column layout
        assessment_table = doc.add_table(rows=1, cols=3)
        assessment_table.style = 'Table Grid'
        assessment_table.autofit = False

        # Define assessment items with icons
        assessments = [
            ('formative_assessment', 'Formative', LIGHT_BLUE),
            ('summative_assessment', 'Summative', SOFT_GREEN),
            ('weekly_deliverable', 'Deliverable', CREAM_YELLOW),
        ]

        for idx, (key, label, color) in enumerate(assessments):
            if week_data.get(key):
                cell = assessment_table.rows[0].cells[idx]
                cell.width = Inches(2.2)

                # Background color
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}" w:val="clear"/>')
                cell._tc.get_or_add_tcPr().append(shading)

                # Label
                p = cell.paragraphs[0]
                run = p.add_run(label)
                run.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = NAVY_BLUE
                p.paragraph_format.space_after = Pt(6)

                # Content
                p = cell.add_paragraph()
                run = p.add_run(week_data[key])
                run.font.size = Pt(10)
                run.font.color.rgb = DARK_GRAY

        doc.add_paragraph()

    # === DAILY SECTIONS ===
    days = week_data.get('days', [])
    day_names = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']

    for i, day in enumerate(days, 1):
        doc.add_page_break()

        # Day header - Tab-style banner with topic
        day_name = day.get('day_label') or (day_names[i-1] if i <= len(day_names) else f"Day {i}")

        # Create tab-style header with day badge + topic bar
        day_header_table = doc.add_table(rows=1, cols=2)
        day_header_table.style = 'Table Grid'
        day_header_table.autofit = False

        # Day number "tab"
        tab_cell = day_header_table.rows[0].cells[0]
        tab_cell.width = Inches(1.2)
        tab_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{ACCENT_BLUE}" w:val="clear"/>')
        tab_cell._tc.get_or_add_tcPr().append(tab_shading)

        p = tab_cell.paragraphs[0]
        run = p.add_run(f"DAY {i}")
        run.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = 'Cambria'
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        p = tab_cell.add_paragraph()
        run = p.add_run(day_name)
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Topic bar
        topic_cell = day_header_table.rows[0].cells[1]
        topic_cell.width = Inches(5.7)
        topic_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
        topic_cell._tc.get_or_add_tcPr().append(topic_shading)

        p = topic_cell.paragraphs[0]
        run = p.add_run(day.get('topic', 'Untitled'))
        run.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = 'Cambria'
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        # Vertically center - left align
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT

        doc.add_paragraph()  # Spacing after header

        # Learning Objectives (compact card style)
        if day.get('objectives'):
            add_section_header("Learning Objectives", level=2)

            obj_box = doc.add_table(rows=1, cols=1)
            obj_box.style = 'Table Grid'
            obj_cell = obj_box.rows[0].cells[0]
            obj_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_BLUE}" w:val="clear"/>')
            obj_cell._tc.get_or_add_tcPr().append(obj_shading)

            for idx, obj in enumerate(day['objectives']):
                p = obj_cell.paragraphs[0] if idx == 0 else obj_cell.add_paragraph()
                run = p.add_run(obj)
                run.font.color.rgb = DARK_GRAY
                run.font.size = Pt(10)
                p.paragraph_format.space_after = Pt(2)

        # Materials (inline badges style)
        if day.get('day_materials'):
            add_section_header("Materials", level=2)

            p = doc.add_paragraph()
            for idx, mat in enumerate(day['day_materials']):
                if idx > 0:
                    p.add_run("  •  ").font.color.rgb = MEDIUM_GRAY
                run = p.add_run(mat)
                run.font.color.rgb = DARK_GRAY
                run.font.size = Pt(10)

        # Schedule Table (90 minutes) - Enhanced with highlighted time column
        if day.get('schedule') or day.get('activities'):
            add_section_header("Schedule (90 minutes)", level=2)

            activities = day.get('schedule') or day.get('activities', [])
            if activities:
                schedule_table = doc.add_table(rows=1, cols=3)
                schedule_table.style = 'Table Grid'
                schedule_table.autofit = False

                # Header row with enhanced styling
                header_cells = schedule_table.rows[0].cells
                header_texts = ["Time", "Activity", "Details"]
                for idx, header_text in enumerate(header_texts):
                    run = header_cells[idx].paragraphs[0].add_run(header_text)
                    run.bold = True
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    # Navy header background
                    h_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
                    header_cells[idx]._tc.get_or_add_tcPr().append(h_shading)

                # Set column widths
                header_cells[0].width = Inches(0.9)
                header_cells[1].width = Inches(1.8)
                header_cells[2].width = Inches(4.2)

                # Add activity rows with highlighted time column
                for row_idx, activity in enumerate(activities):
                    row = schedule_table.add_row()
                    cells = row.cells

                    if isinstance(activity, dict):
                        # Time cell - always highlighted with accent color
                        time_text = activity.get('time', activity.get('duration', ''))
                        p = cells[0].paragraphs[0]
                        run = p.add_run(time_text)
                        run.bold = True
                        run.font.size = Pt(10)
                        run.font.color.rgb = NAVY_BLUE
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        time_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_BLUE}" w:val="clear"/>')
                        cells[0]._tc.get_or_add_tcPr().append(time_shading)

                        # Activity name - bold
                        p = cells[1].paragraphs[0]
                        run = p.add_run(activity.get('name', activity.get('activity', '')))
                        run.bold = True
                        run.font.size = Pt(10)
                        run.font.color.rgb = DARK_GRAY

                        # Description
                        p = cells[2].paragraphs[0]
                        run = p.add_run(activity.get('description', ''))
                        run.font.size = Pt(10)
                        run.font.color.rgb = DARK_GRAY
                    else:
                        cells[1].text = str(activity)

                    # Alternating row colors for non-time cells
                    if row_idx % 2 == 1:
                        for cell in cells[1:]:  # Skip time column
                            row_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_GRAY}" w:val="clear"/>')
                            cell._tc.get_or_add_tcPr().append(row_shading)

            doc.add_paragraph()

        # Key Vocabulary - Card-style two-column layout
        if day.get('vocabulary'):
            add_section_header("Key Vocabulary", level=2)

            vocab_items = list(day['vocabulary'].items())
            vocab_table = doc.add_table(rows=0, cols=2)
            vocab_table.style = 'Table Grid'
            vocab_table.autofit = False

            # Create vocabulary cards in two columns
            for i in range(0, len(vocab_items), 2):
                row = vocab_table.add_row()
                for j in range(2):
                    if i + j < len(vocab_items):
                        term, definition = vocab_items[i + j]
                        cell = row.cells[j]
                        cell.width = Inches(3.4)

                        # Card background
                        card_bg = LIGHT_BLUE if (i // 2) % 2 == 0 else LIGHT_GRAY
                        v_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{card_bg}" w:val="clear"/>')
                        cell._tc.get_or_add_tcPr().append(v_shading)

                        # Term (bold, navy)
                        p = cell.paragraphs[0]
                        run = p.add_run(term)
                        run.bold = True
                        run.font.size = Pt(11)
                        run.font.color.rgb = NAVY_BLUE
                        p.paragraph_format.space_after = Pt(4)

                        # Definition
                        p = cell.add_paragraph()
                        run = p.add_run(definition)
                        run.font.size = Pt(10)
                        run.font.color.rgb = DARK_GRAY

            doc.add_paragraph()

        # Differentiation Strategies - Three-column card layout
        if day.get('differentiation'):
            add_section_header("Differentiation Strategies", level=2)

            diff = day['differentiation']
            if isinstance(diff, dict):
                diff_table = doc.add_table(rows=1, cols=len(diff))
                diff_table.style = 'Table Grid'
                diff_table.autofit = False

                # Color coding for differentiation levels
                diff_colors = {
                    'Struggling': "FFEBEE",  # Light red
                    'On-Level': LIGHT_BLUE,  # Light blue
                    'Advanced': SOFT_GREEN,  # Light green
                    'ELL': "FFF3E0",  # Light orange
                }

                for idx, (level, strategy) in enumerate(diff.items()):
                    cell = diff_table.rows[0].cells[idx]
                    cell.width = Inches(6.8 / len(diff))

                    # Get color for level
                    bg_color = diff_colors.get(level, LIGHT_GRAY)
                    d_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>')
                    cell._tc.get_or_add_tcPr().append(d_shading)

                    # Level label
                    p = cell.paragraphs[0]
                    run = p.add_run(level)
                    run.bold = True
                    run.font.size = Pt(11)
                    run.font.color.rgb = NAVY_BLUE
                    p.paragraph_format.space_after = Pt(4)

                    # Strategy
                    p = cell.add_paragraph()
                    run = p.add_run(strategy)
                    run.font.size = Pt(10)
                    run.font.color.rgb = DARK_GRAY
            else:
                p = doc.add_paragraph()
                run = p.add_run(diff)
                run.font.color.rgb = DARK_GRAY

            doc.add_paragraph()

        # Teacher Notes (Sticky-note style box)
        if day.get('teacher_notes'):
            add_section_header("Teacher Notes", level=2)

            # Create sticky-note style box with slight rotation effect
            note_table = doc.add_table(rows=1, cols=2)
            note_table.style = 'Table Grid'
            note_table.autofit = False

            # Left accent bar (like a post-it tab)
            tab_cell = note_table.rows[0].cells[0]
            tab_cell.width = Inches(0.15)
            tab_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="FFD93D" w:val="clear"/>')  # Yellow tab
            tab_cell._tc.get_or_add_tcPr().append(tab_shading)

            # Note content
            note_cell = note_table.rows[0].cells[1]
            note_cell.width = Inches(6.65)
            note_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM_YELLOW}" w:val="clear"/>')
            note_cell._tc.get_or_add_tcPr().append(note_shading)

            p = note_cell.paragraphs[0]
            run = p.add_run(day['teacher_notes'])
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(10)
            run.font.italic = True

            doc.add_paragraph()

    # === END OF WEEK SECTIONS ===

    # Week Vocabulary Summary
    if week_data.get('vocabulary_summary'):
        doc.add_page_break()
        add_section_header("Week Vocabulary Summary", level=1)

        # Card layout for vocabulary categories
        for category, terms in week_data['vocabulary_summary'].items():
            cat_table = doc.add_table(rows=1, cols=1)
            cat_table.style = 'Table Grid'
            cell = cat_table.rows[0].cells[0]
            cat_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_BLUE}" w:val="clear"/>')
            cell._tc.get_or_add_tcPr().append(cat_shading)

            p = cell.paragraphs[0]
            run = p.add_run(category)
            run.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = NAVY_BLUE
            p.paragraph_format.space_after = Pt(4)

            p = cell.add_paragraph()
            run = p.add_run(terms)
            run.font.size = Pt(10)
            run.font.color.rgb = DARK_GRAY

            doc.add_paragraph()

    # Teacher Notes (end of week) - Pull-quote box style
    if week_data.get('teacher_notes'):
        add_section_header("Teacher Notes", level=1)

        notes_table = doc.add_table(rows=1, cols=2)
        notes_table.style = 'Table Grid'
        notes_table.autofit = False

        # Yellow accent bar
        accent_cell = notes_table.rows[0].cells[0]
        accent_cell.width = Inches(0.15)
        accent_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="FFD93D" w:val="clear"/>')
        accent_cell._tc.get_or_add_tcPr().append(accent_shading)

        # Notes content
        notes_cell = notes_table.rows[0].cells[1]
        notes_cell.width = Inches(6.65)
        notes_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM_YELLOW}" w:val="clear"/>')
        notes_cell._tc.get_or_add_tcPr().append(notes_shading)

        for idx, note in enumerate(week_data['teacher_notes']):
            p = notes_cell.paragraphs[0] if idx == 0 else notes_cell.add_paragraph()
            run = p.add_run(f"• {note}")
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(10)
            run.font.italic = True

        doc.add_paragraph()

    # Standards Alignment - Professional footer style
    if week_data.get('standards_alignment'):
        add_section_header("Standards Alignment", level=1)

        std_table = doc.add_table(rows=1, cols=1)
        std_table.style = 'Table Grid'
        cell = std_table.rows[0].cells[0]
        std_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_GRAY}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(std_shading)

        p = cell.paragraphs[0]
        run = p.add_run(week_data['standards_alignment'])
        run.font.color.rgb = MEDIUM_GRAY
        run.font.size = Pt(9)
        run.font.italic = True

    # Save document in week folder
    week_folder = get_week_folder(week_num)
    unit_slug = sanitize_slug(unit_name, max_length=20) if unit_name else 'Lessons'
    filename = f"Week{week_num}_{unit_slug}_TeacherHandout.docx"
    output_path = os.path.join(week_folder, filename)

    doc.save(output_path)
    return output_path


def generate_student_handout(handout_data, week_num, handout_name):
    """Generate a Canva-quality student handout with enhanced visual design."""
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    doc = Document()

    # Use shared color palette
    NAVY_BLUE = DOC_NAVY_BLUE
    DARK_GRAY = DOC_DARK_GRAY
    MEDIUM_GRAY = DOC_MEDIUM_GRAY
    LIGHT_BLUE = DOC_LIGHT_BLUE
    LIGHT_GRAY = "F8F9FA"  # Slightly lighter for student handouts
    ACCENT_BLUE = DOC_ACCENT_BLUE
    CREAM_YELLOW = DOC_CREAM_YELLOW
    SOFT_GREEN = DOC_SOFT_GREEN

    setup_document(doc, top=0.7, bottom=0.7, left=0.8, right=0.8,
                   line_spacing=1.25, space_after=8)

    def add_section_header(text):
        return add_doc_section_header(doc, text, level=1)

    # === HEADER BANNER ===
    add_doc_header_banner(doc, handout_data.get('title', 'Student Handout'),
                          subtitle_text=handout_data.get('subtitle'))

    doc.add_paragraph()  # Extra spacing

    # === INSTRUCTIONS BOX (Card style) ===
    if handout_data.get('instructions'):
        add_section_header("Instructions")

        inst_table = doc.add_table(rows=1, cols=1)
        inst_table.style = 'Table Grid'
        cell = inst_table.rows[0].cells[0]
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_BLUE}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)

        p = cell.paragraphs[0]
        run = p.add_run(handout_data['instructions'])
        run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(11)

        doc.add_paragraph()

    # === MAIN CONTENT SECTIONS ===
    for section in handout_data.get('sections', []):
        heading = section.get('heading', '')
        add_section_header(heading if heading else "Content")

        if section.get('content'):
            p = doc.add_paragraph()
            run = p.add_run(section['content'])
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(11)
            p.paragraph_format.space_after = Pt(10)

        if section.get('items'):
            # Use circular badge numbered items or clean bullets
            if section.get('numbered'):
                items_table = doc.add_table(rows=0, cols=2)
                items_table.style = 'Table Grid'
                items_table.autofit = False

                for idx, item in enumerate(section['items'], 1):
                    row = items_table.add_row()

                    # Circular number badge cell
                    num_cell = row.cells[0]
                    num_cell.width = Inches(0.45)
                    p = num_cell.paragraphs[0]
                    run = p.add_run(str(idx))
                    run.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.size = Pt(12)
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    # Navy background for circular effect
                    badge_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
                    num_cell._tc.get_or_add_tcPr().append(badge_shading)

                    # Content cell with alternating subtle background
                    content_cell = row.cells[1]
                    content_cell.width = Inches(6.05)
                    p = content_cell.paragraphs[0]
                    run = p.add_run(item)
                    run.font.color.rgb = DARK_GRAY
                    run.font.size = Pt(11)

                    if idx % 2 == 0:
                        row_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_GRAY}" w:val="clear"/>')
                        content_cell._tc.get_or_add_tcPr().append(row_shading)
            else:
                for item in section['items']:
                    p = doc.add_paragraph()
                    run = p.add_run(f"• {item}")
                    run.font.color.rgb = DARK_GRAY
                    run.font.size = Pt(11)
                    p.paragraph_format.left_indent = Inches(0.25)

        # Blank lines for writing (styled with subtle gray)
        if section.get('blank_lines'):
            doc.add_paragraph()
            for _ in range(section['blank_lines']):
                p = doc.add_paragraph()
                run = p.add_run('_' * 85)
                run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                p.paragraph_format.space_after = Pt(14)

        doc.add_paragraph()

    # === QUESTIONS SECTION (Enhanced with number badges) ===
    if handout_data.get('questions'):
        add_section_header("Questions")

        for i, q in enumerate(handout_data['questions'], 1):
            q_table = doc.add_table(rows=1, cols=2)
            q_table.style = 'Table Grid'
            q_table.autofit = False

            # Question number badge
            num_cell = q_table.rows[0].cells[0]
            num_cell.width = Inches(0.5)
            badge_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1A3C6E" w:val="clear"/>')
            num_cell._tc.get_or_add_tcPr().append(badge_shading)

            p = num_cell.paragraphs[0]
            run = p.add_run(str(i))
            run.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(14)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Question and answer area
            q_cell = q_table.rows[0].cells[1]
            q_cell.width = Inches(6.0)

            # Alternating background
            if i % 2 == 0:
                q_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{LIGHT_GRAY}" w:val="clear"/>')
                q_cell._tc.get_or_add_tcPr().append(q_shading)

            # Question text
            p = q_cell.paragraphs[0]
            run = p.add_run(q)
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(11)
            p.paragraph_format.space_after = Pt(10)

            # Answer lines
            for _ in range(3):
                p = q_cell.add_paragraph()
                run = p.add_run('_' * 80)
                run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                p.paragraph_format.space_after = Pt(8)

            doc.add_paragraph()

    # === VOCABULARY SECTION (Two-column card layout) ===
    if handout_data.get('vocabulary'):
        add_section_header("Vocabulary")

        vocab_items = list(handout_data['vocabulary'].items())
        vocab_table = doc.add_table(rows=0, cols=2)
        vocab_table.style = 'Table Grid'
        vocab_table.autofit = False

        # Create vocabulary cards in two columns
        for i in range(0, len(vocab_items), 2):
            row = vocab_table.add_row()
            for j in range(2):
                if i + j < len(vocab_items):
                    term, definition = vocab_items[i + j]
                    cell = row.cells[j]
                    cell.width = Inches(3.25)

                    # Card background - alternating colors
                    card_bg = LIGHT_BLUE if (i // 2) % 2 == 0 else LIGHT_GRAY
                    v_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{card_bg}" w:val="clear"/>')
                    cell._tc.get_or_add_tcPr().append(v_shading)

                    # Term (bold, navy)
                    p = cell.paragraphs[0]
                    run = p.add_run(term)
                    run.bold = True
                    run.font.size = Pt(11)
                    run.font.color.rgb = NAVY_BLUE
                    p.paragraph_format.space_after = Pt(4)

                    # Definition
                    p = cell.add_paragraph()
                    run = p.add_run(definition)
                    run.font.size = Pt(10)
                    run.font.color.rgb = DARK_GRAY

        doc.add_paragraph()

    # === TIPS/NOTES SECTION (Pull-quote style with accent) ===
    if handout_data.get('tips') or handout_data.get('notes'):
        add_section_header("Tips & Notes")

        tips_table = doc.add_table(rows=1, cols=2)
        tips_table.style = 'Table Grid'
        tips_table.autofit = False

        # Yellow accent bar (like a highlight strip)
        accent_cell = tips_table.rows[0].cells[0]
        accent_cell.width = Inches(0.12)
        accent_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="FFD93D" w:val="clear"/>')
        accent_cell._tc.get_or_add_tcPr().append(accent_shading)

        # Content cell
        content_cell = tips_table.rows[0].cells[1]
        content_cell.width = Inches(6.38)
        content_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM_YELLOW}" w:val="clear"/>')
        content_cell._tc.get_or_add_tcPr().append(content_shading)

        tips = handout_data.get('tips') or handout_data.get('notes', [])
        if isinstance(tips, list):
            for idx, tip in enumerate(tips):
                p = content_cell.paragraphs[0] if idx == 0 else content_cell.add_paragraph()
                run = p.add_run(f"- {tip}")
                run.font.color.rgb = DARK_GRAY
                run.font.size = Pt(10)
        else:
            p = content_cell.paragraphs[0]
            run = p.add_run(tips)
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(10)

    # Save document in week folder
    week_folder = get_week_folder(week_num)
    name_slug = sanitize_slug(handout_name)
    filename = f"{name_slug}_StudentHandout.docx"
    output_path = os.path.join(week_folder, filename)

    doc.save(output_path)
    return output_path


def generate_bell_ringer_slides(week_data):
    """Generate Bell Ringer slides as PowerPoint for Canva upload."""
    week_num = week_data.get('week', '')
    unit_name = week_data.get('unit', '')
    days = week_data.get('days', [])
    day_names = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']

    # Colors
    NAVY_BLUE = PptxRGBColor(0x1a, 0x3c, 0x6e)
    WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BLUE = PptxRGBColor(0xD6, 0xE3, 0xF8)

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    slides_created = []

    for i, day in enumerate(days, 1):
        # Use day_label if provided, otherwise fall back to default day names
        day_name = day.get('day_label') or (day_names[i-1] if i <= len(day_names) else f"Day {i}")
        topic = day.get('topic', '')

        # Find Bell Ringer in schedule
        bell_ringer_text = None
        schedule = day.get('schedule', [])
        for activity in schedule:
            if isinstance(activity, dict):
                activity_name = activity.get('name', activity.get('activity', '')).lower()
                if 'bell ringer' in activity_name or 'bellringer' in activity_name or 'warm up' in activity_name or 'warmup' in activity_name:
                    bell_ringer_text = activity.get('description', '')
                    break

        # Use placeholder if no Bell Ringer found
        if not bell_ringer_text:
            bell_ringer_text = "[Add Bell Ringer prompt]"

        # Create blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add navy background rectangle (full slide)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(0), PptxInches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = NAVY_BLUE
        bg_shape.line.fill.background()

        # Add decorative light blue accent bar at top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(0), PptxInches(0),
            prs.slide_width, PptxInches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = LIGHT_BLUE
        accent_bar.line.fill.background()

        # Add "BELL RINGER" title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.8),
            PptxInches(12.333), PptxInches(1.0)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.add_run()
        title_run.text = "BELL RINGER"
        title_run.font.name = "Cambria"
        title_run.font.size = PptxPt(54)
        title_run.font.bold = True
        title_run.font.color.rgb = WHITE

        # Add day info subtitle
        day_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(1.8),
            PptxInches(12.333), PptxInches(0.6)
        )
        day_frame = day_box.text_frame
        day_frame.word_wrap = True
        day_para = day_frame.paragraphs[0]
        day_para.alignment = PP_ALIGN.CENTER
        day_run = day_para.add_run()
        day_run.text = f"Week {week_num} • {day_name}"
        day_run.font.name = "Calibri"
        day_run.font.size = PptxPt(24)
        day_run.font.color.rgb = LIGHT_BLUE

        # Add content box with light background
        content_box_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PptxInches(0.75), PptxInches(2.8),
            PptxInches(11.833), PptxInches(3.8)
        )
        content_box_shape.fill.solid()
        content_box_shape.fill.fore_color.rgb = WHITE
        content_box_shape.line.fill.background()

        # Add bell ringer prompt text
        prompt_box = slide.shapes.add_textbox(
            PptxInches(1.25), PptxInches(3.2),
            PptxInches(10.833), PptxInches(3.0)
        )
        prompt_frame = prompt_box.text_frame
        prompt_frame.word_wrap = True
        prompt_frame.auto_size = None

        prompt_para = prompt_frame.paragraphs[0]
        prompt_para.alignment = PP_ALIGN.CENTER
        prompt_run = prompt_para.add_run()
        prompt_run.text = bell_ringer_text
        prompt_run.font.name = "Calibri"
        prompt_run.font.size = PptxPt(32)
        prompt_run.font.color.rgb = NAVY_BLUE

        # Center text vertically
        prompt_frame.paragraphs[0].space_before = PptxPt(20)

        slides_created.append({
            'day': i,
            'day_name': day_name,
            'prompt': bell_ringer_text
        })

    # Only save if we have slides
    if not slides_created:
        return None, []

    # Save presentation in week folder
    week_folder = get_week_folder(week_num)
    filename = f"Week{week_num}_BellRinger_Slides.pptx"
    output_path = os.path.join(week_folder, filename)

    prs.save(output_path)
    return output_path, slides_created


def generate_daily_presentation(day_data, week_num, day_num, unit_name=''):
    """
    Generate a full 90-minute lesson presentation for a single day.

    Structure:
    1. Bell Ringer - Question/prompt with background image
    2. Agenda - Visual timeline of day's activities
    3. Direct Instruction - Content slides with images
    4. Guided Practice - Activity instructions
    5. Hands-On Activity - Steps and expectations
    6. Wrap-Up - Key takeaways and exit ticket
    """
    day_names = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']
    day_name = day_data.get('day_label') or (day_names[day_num-1] if day_num <= len(day_names) else f"Day {day_num}")
    topic = day_data.get('topic', 'Lesson')

    # Get color theme for unit
    colors = UNIT_COLOR_THEMES.get(unit_name, DEFAULT_COLOR_THEME)
    PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR = colors
    WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
    DARK_GRAY = PptxRGBColor(0x33, 0x33, 0x33)

    # Track media for logging
    media_log = {'images': [], 'videos': []}

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_background(slide, color):
        """Add solid color background to slide."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        # Send to back
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def add_title_bar(slide, title_text, subtitle_text=None):
        """Add a colored title bar at top of slide."""
        # Title bar background
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptxInches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_COLOR
        bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.25), PptxInches(12.333), PptxInches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name = "Cambria"
        run.font.size = PptxPt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            run2 = p2.add_run()
            run2.text = subtitle_text
            run2.font.name = "Calibri"
            run2.font.size = PptxPt(18)
            run2.font.color.rgb = SECONDARY_COLOR

    def add_content_with_image(slide, title, bullets, image_query):
        """Add content slide with bullets on left and image on right."""
        add_title_bar(slide, title)

        # Bullet content on left
        content_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(6.5), PptxInches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.name = "Calibri"
            run.font.size = PptxPt(24)
            run.font.color.rgb = DARK_GRAY
            p.space_after = PptxPt(12)

        # Try to add image on right
        if image_query:
            image_data, image_url = get_topic_image(image_query)
            if image_data:
                try:
                    slide.shapes.add_picture(image_data, PptxInches(7.2), PptxInches(1.5), width=PptxInches(5.5))
                    media_log['images'].append({'query': image_query, 'url': image_url})
                except Exception as e:
                    print(f"Could not add image: {e}", file=sys.stderr)

    # =========================================================================
    # SLIDE 1: BELL RINGER
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY_COLOR)

    # Find Bell Ringer text
    bell_ringer_text = "[Add Bell Ringer prompt]"
    for activity in day_data.get('schedule', []):
        if isinstance(activity, dict):
            name = activity.get('name', '').lower()
            if 'bell ringer' in name or 'warm up' in name:
                bell_ringer_text = activity.get('description', bell_ringer_text)
                break

    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptxInches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    # "BELL RINGER" title
    title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.8), PptxInches(12.333), PptxInches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "BELL RINGER"
    run.font.name = "Cambria"
    run.font.size = PptxPt(54)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Day/Week info
    info_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.7), PptxInches(12.333), PptxInches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Week {week_num} • {day_name}"
    run.font.name = "Calibri"
    run.font.size = PptxPt(22)
    run.font.color.rgb = SECONDARY_COLOR

    # Content box
    content_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.75), PptxInches(2.6), PptxInches(11.833), PptxInches(4.2))
    content_shape.fill.solid()
    content_shape.fill.fore_color.rgb = WHITE
    content_shape.line.fill.background()

    # Bell ringer prompt
    prompt_box = slide.shapes.add_textbox(PptxInches(1.2), PptxInches(3.2), PptxInches(10.9), PptxInches(3.2))
    tf = prompt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = bell_ringer_text
    run.font.name = "Calibri"
    run.font.size = PptxPt(32)
    run.font.color.rgb = PRIMARY_COLOR

    # =========================================================================
    # SLIDE 2: AGENDA
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SECONDARY_COLOR)
    add_title_bar(slide, "TODAY'S AGENDA", topic)

    schedule = day_data.get('schedule', [])
    y_pos = 1.6
    for activity in schedule:
        if isinstance(activity, dict):
            time = activity.get('time', '')
            name = activity.get('name', '')
            if time and name:
                # Time badge
                time_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.5), PptxInches(y_pos), PptxInches(1.3), PptxInches(0.5))
                time_box.fill.solid()
                time_box.fill.fore_color.rgb = PRIMARY_COLOR
                time_box.line.fill.background()

                # Add time text to shape
                time_tf = time_box.text_frame
                time_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                time_run = time_tf.paragraphs[0].add_run()
                time_run.text = time
                time_run.font.name = "Calibri"
                time_run.font.size = PptxPt(14)
                time_run.font.bold = True
                time_run.font.color.rgb = WHITE

                # Activity name
                name_box = slide.shapes.add_textbox(PptxInches(2.0), PptxInches(y_pos + 0.08), PptxInches(10), PptxInches(0.5))
                name_tf = name_box.text_frame
                name_run = name_tf.paragraphs[0].add_run()
                name_run.text = name
                name_run.font.name = "Calibri"
                name_run.font.size = PptxPt(22)
                name_run.font.color.rgb = DARK_GRAY

                y_pos += 0.7

    # =========================================================================
    # SLIDES 3+: DIRECT INSTRUCTION CONTENT
    # =========================================================================
    objectives = day_data.get('objectives', [])
    vocabulary = day_data.get('vocabulary', {})

    # Objectives slide
    if objectives:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, SECONDARY_COLOR)
        add_content_with_image(slide, "LEARNING OBJECTIVES", objectives, topic)

    # Vocabulary slide (if exists)
    if vocabulary:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, SECONDARY_COLOR)
        add_title_bar(slide, "KEY VOCABULARY")

        y_pos = 1.5
        for term, definition in vocabulary.items():
            # Term box
            term_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.5), PptxInches(y_pos), PptxInches(3), PptxInches(0.6))
            term_box.fill.solid()
            term_box.fill.fore_color.rgb = PRIMARY_COLOR
            term_box.line.fill.background()

            term_tf = term_box.text_frame
            term_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            term_tf.paragraphs[0].space_before = PptxPt(8)
            term_run = term_tf.paragraphs[0].add_run()
            term_run.text = term
            term_run.font.name = "Calibri"
            term_run.font.size = PptxPt(18)
            term_run.font.bold = True
            term_run.font.color.rgb = WHITE

            # Definition
            def_box = slide.shapes.add_textbox(PptxInches(3.7), PptxInches(y_pos + 0.1), PptxInches(9), PptxInches(0.5))
            def_tf = def_box.text_frame
            def_run = def_tf.paragraphs[0].add_run()
            def_run.text = definition
            def_run.font.name = "Calibri"
            def_run.font.size = PptxPt(18)
            def_run.font.color.rgb = DARK_GRAY

            y_pos += 0.8

    # Topic content slide with image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SECONDARY_COLOR)
    topic_bullets = [f"Today's focus: {topic}"]
    if objectives:
        topic_bullets.extend(objectives[:3])
    add_content_with_image(slide, topic.upper(), topic_bullets, topic)

    # =========================================================================
    # VIDEO SLIDE (if relevant video found)
    # =========================================================================
    video_url, video_title = search_youtube_video(topic)
    if video_url:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, SECONDARY_COLOR)
        add_title_bar(slide, "VIDEO", video_title[:50] + "..." if len(video_title) > 50 else video_title)

        # Video placeholder box
        video_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(1.5), PptxInches(1.8), PptxInches(10.333), PptxInches(5.2))
        video_box.fill.solid()
        video_box.fill.fore_color.rgb = PptxRGBColor(0x20, 0x20, 0x20)
        video_box.line.fill.background()

        # Play button icon (triangle)
        play_btn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, PptxInches(6.1), PptxInches(3.8), PptxInches(1.2), PptxInches(1.2))
        play_btn.fill.solid()
        play_btn.fill.fore_color.rgb = WHITE
        play_btn.rotation = 90

        # Video URL text
        url_box = slide.shapes.add_textbox(PptxInches(1.5), PptxInches(6.2), PptxInches(10.333), PptxInches(0.5))
        url_tf = url_box.text_frame
        url_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        url_run = url_tf.paragraphs[0].add_run()
        url_run.text = video_url
        url_run.font.name = "Calibri"
        url_run.font.size = PptxPt(12)
        url_run.font.color.rgb = ACCENT_COLOR

        media_log['videos'].append({'title': video_title, 'url': video_url})

    # =========================================================================
    # ACTIVITY/PRACTICE SLIDES
    # =========================================================================
    # Find guided practice and hands-on activities from schedule
    for activity in schedule:
        if isinstance(activity, dict):
            name = activity.get('name', '').lower()
            desc = activity.get('description', '')

            if any(word in name for word in ['practice', 'activity', 'hands-on', 'work time', 'project']):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_background(slide, SECONDARY_COLOR)
                add_title_bar(slide, activity.get('name', 'ACTIVITY').upper())

                # Activity description
                desc_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(12.333), PptxInches(5.5))
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = desc
                run.font.name = "Calibri"
                run.font.size = PptxPt(28)
                run.font.color.rgb = DARK_GRAY

    # =========================================================================
    # WRAP-UP SLIDE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(12.333), PptxInches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "WRAP-UP"
    run.font.name = "Cambria"
    run.font.size = PptxPt(48)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Key takeaways box
    takeaway_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.5), PptxInches(1.8), PptxInches(6), PptxInches(4.5))
    takeaway_shape.fill.solid()
    takeaway_shape.fill.fore_color.rgb = WHITE
    takeaway_shape.line.fill.background()

    takeaway_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(2.0), PptxInches(5.5), PptxInches(4.0))
    tf = takeaway_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Takeaways:"
    run.font.name = "Cambria"
    run.font.size = PptxPt(22)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_COLOR

    if objectives:
        for obj in objectives[:3]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"- {obj[:60]}..." if len(obj) > 60 else f"- {obj}"
            run.font.name = "Calibri"
            run.font.size = PptxPt(16)
            run.font.color.rgb = DARK_GRAY

    # Exit ticket box
    exit_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(6.833), PptxInches(1.8), PptxInches(6), PptxInches(4.5))
    exit_shape.fill.solid()
    exit_shape.fill.fore_color.rgb = ACCENT_COLOR
    exit_shape.line.fill.background()

    exit_box = slide.shapes.add_textbox(PptxInches(7.1), PptxInches(2.0), PptxInches(5.5), PptxInches(4.0))
    tf = exit_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exit Ticket"
    run.font.name = "Cambria"
    run.font.size = PptxPt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Find exit ticket from schedule
    exit_text = "What did you learn today?"
    for activity in schedule:
        if isinstance(activity, dict):
            name = activity.get('name', '').lower()
            if 'wrap' in name or 'exit' in name or 'reflection' in name:
                exit_text = activity.get('description', exit_text)
                break

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = exit_text
    run.font.name = "Calibri"
    run.font.size = PptxPt(18)
    run.font.color.rgb = WHITE

    # Save presentation
    week_folder = get_week_folder(week_num)
    topic_slug = sanitize_slug(topic)
    filename = f"Day{day_num}_{topic_slug}_Presentation.pptx"
    output_path = os.path.join(week_folder, filename)

    prs.save(output_path)
    return output_path, media_log


def generate_week(data):
    """Generate all documents for a week: CTE plans, teacher handout, student handouts, and daily presentations."""
    results = {
        'cte_plans': [],
        'teacher_handout': None,
        'student_handouts': [],
        'daily_presentations': [],
        'media_log': [],
        'week_folder': None
    }

    week_num = data.get('week', '1')
    unit_name = data.get('unit', '')
    days = data.get('days', [])
    skip_presentations = data.get('skip_presentations', False)

    # Store the week folder path
    results['week_folder'] = get_week_folder(week_num)

    # Generate individual CTE lesson plans
    for i, day in enumerate(days, 1):
        path = generate_cte_lesson_plan(day, week_num, i)
        results['cte_plans'].append(path)

    # Generate teacher handout
    results['teacher_handout'] = generate_teacher_handout(data)

    # Generate student handouts
    for handout in data.get('student_handouts', []):
        path = generate_student_handout(handout, week_num, handout.get('name', 'Handout'))
        results['student_handouts'].append(path)

    # Generate daily lesson presentations (unless skip_presentations is True)
    all_media_log = []
    if not skip_presentations:
        for i, day in enumerate(days, 1):
            try:
                pres_path, media_log = generate_daily_presentation(day, week_num, i, unit_name)
                results['daily_presentations'].append(pres_path)
                all_media_log.extend(media_log)
            except Exception as e:
                print(f"Warning: Could not generate presentation for Day {i}: {e}", file=sys.stderr)

        # Write media log file
        if all_media_log:
            media_log_path = os.path.join(results['week_folder'], f"Week{week_num}_Media_Log.txt")
            with open(media_log_path, 'w') as f:
                f.write(f"Media Log - Week {week_num}: {unit_name}\n")
                f.write("=" * 60 + "\n\n")
                for entry in all_media_log:
                    f.write(f"{entry}\n")
            results['media_log'] = media_log_path

    return results


def validate_lesson_data(data):
    """Validate the structure of lesson plan JSON data."""
    if not isinstance(data, dict):
        raise ValueError("Input must be a JSON object")

    if 'days' in data:
        if not isinstance(data['days'], list) or not data['days']:
            raise ValueError("'days' must be a non-empty array")
        for i, day in enumerate(data['days'], 1):
            if not isinstance(day, dict):
                raise ValueError(f"Day {i} must be a JSON object")
            if not day.get('topic'):
                raise ValueError(f"Day {i} is missing required 'topic' field")
    else:
        if not data.get('topic'):
            raise ValueError("Single lesson is missing required 'topic' field")


def print_results(results):
    """Print generation results summary."""
    print("SUCCESS: Weekly lesson plans generated")
    print(f"Week Folder: {results['week_folder']}")
    print(f"CTE Plans: {len(results['cte_plans'])}")
    for path in results['cte_plans']:
        print(f"  - {os.path.basename(path)}")
    print(f"Teacher Handout: {os.path.basename(results['teacher_handout'])}")
    for path in results['student_handouts']:
        print(f"Student Handout: {os.path.basename(path)}")
    if results['daily_presentations']:
        print(f"Daily Presentations: {len(results['daily_presentations'])}")
        for path in results['daily_presentations']:
            print(f"  - {os.path.basename(path)}")
    if results.get('media_log'):
        print(f"Media Log: {os.path.basename(results['media_log'])}")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python generate-lesson-plan.py '<json_data>'", file=sys.stderr)
        sys.exit(1)

    try:
        data = json.loads(sys.argv[1])
        validate_lesson_data(data)

        if 'days' in data:
            results = generate_week(data)
            print_results(results)
        else:
            week = data.get('week', '1')
            output_path = generate_cte_lesson_plan(data, week, 1)
            print(f"SUCCESS: {output_path}")

    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON - {e}", file=sys.stderr)
        sys.exit(1)
    except ValueError as e:
        print(f"ERROR: Invalid input - {e}", file=sys.stderr)
        sys.exit(1)
    except FileNotFoundError as e:
        print(f"ERROR: File not found - {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        import traceback
        print(f"ERROR: {e}", file=sys.stderr)
        traceback.print_exc()
        sys.exit(1)
